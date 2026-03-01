import subprocess
import os
import webbrowser
from datetime import datetime
from reportlab.pdfgen import canvas
from pptx import Presentation
import psutil
import pyautogui
from ctypes import cast, POINTER
from comtypes import CLSCTX_ALL
from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume

def get_current_time():
    """
    Returns the current local system time and date.
    Use this when the user asks for the time, date, or what day it is.
    """
    now = datetime.now()
    return f"The current date and time is: {now.strftime('%A, %B %d, %Y at %I:%M %p')}"

def open_application(app_name: str):
    """
    Opens a standard Windows application by name.
    Args:
        app_name: The name of the application to open (e.g., 'notepad', 'calc', 'explorer', 'chrome').
    """
    app_name = app_name.lower().strip()
    
    app_map = {
        "notepad": "notepad.exe",
        "calculator": "calc.exe",
        "calc": "calc.exe",
        "file explorer": "explorer.exe",
        "explorer": "explorer.exe",
        "command prompt": "cmd.exe",
        "cmd": "cmd.exe",
        "browser": "start msedge",
        "chrome": "start chrome",
    }
    
    exe_command = app_map.get(app_name, f"start {app_name}")
    
    try:
        subprocess.Popen(exe_command, shell=True)
        return f"Successfully opened {app_name}."
    except Exception as e:
        return f"Failed to open {app_name}. Error: {e}"

def close_application(app_name: str):
    """
    Closes a running Windows application by forcibly terminating its process.
    Args:
        app_name: The name of the application to close (e.g., 'notepad', 'calc', 'chrome').
    """
    app_name = app_name.lower().strip()
    
    # Common mappings
    app_map = {
        "notepad": "notepad.exe",
        "calculator": "calc.exe",
        "calc": "calc.exe",
        "file explorer": "explorer.exe",
        "explorer": "explorer.exe",
        "command prompt": "cmd.exe",
        "cmd": "cmd.exe",
        "browser": "msedge.exe",
        "chrome": "chrome.exe",
    }
    
    target_exe = app_map.get(app_name, f"{app_name}.exe")
    terminated_count = 0
    
    try:
        for proc in psutil.process_iter(['name']):
            try:
                if proc.info['name'] and proc.info['name'].lower() == target_exe:
                    proc.kill()
                    terminated_count += 1
            except (psutil.NoSuchProcess, psutil.AccessDenied, psutil.ZombieProcess):
                pass
                
        if terminated_count > 0:
            return f"Successfully closed {app_name} ({terminated_count} processes terminated)."
        else:
            return f"Could not find any running processes matching '{app_name}'."
    except Exception as e:
        return f"Failed to close {app_name}. Error: {e}"

def open_website(url: str):
    """
    Opens a specific website URL in the user's default web browser.
    Args:
        url: The full URL to open (e.g., 'https://www.youtube.com' or 'https://www.google.com/search?q=puppies').
    """
    try:
        # Ensure the URL has a scheme
        if not url.startswith('http://') and not url.startswith('https://'):
            url = 'https://' + url
            
        webbrowser.open(url)
        return f"Successfully opened the website: {url}"
    except Exception as e:
        return f"Failed to open website {url}. Error: {e}"

def create_file(filename: str, content: str):
    """
    Creates a new text-based file on the user's Desktop with the specified content. 
    Use this if the user asks you to write a note, save a file, or create a document.
    Args:
        filename: The name of the file to create, including the extension (e.g., 'groceries.txt', 'notes.md').
        content: The text content to write inside the file.
    """
    try:
        # Resolve the path to the user's NOVA workspace folder
        nova_path = r"d:\projects\degree\jarvis\voice_ai_assistant_free_ai\NOVA"
        os.makedirs(nova_path, exist_ok=True)
        file_path = os.path.join(nova_path, filename)
        
        # Write the file
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(content)
            
        return f"Successfully created the file '{filename}' in the NOVA workspace."
    except Exception as e:
        return f"Failed to create file '{filename}'. Error: {e}"

def create_pdf(filename: str, content: str):
    """
    Creates a new PDF document on the user's Desktop with the specified content.
    Args:
        filename: The name of the PDF file to create (MUST end with .pdf).
        content: The text content to write inside the PDF.
    """
    try:
        if not filename.endswith('.pdf'):
            filename += '.pdf'
            
        nova_path = r"d:\projects\degree\jarvis\voice_ai_assistant_free_ai\NOVA"
        os.makedirs(nova_path, exist_ok=True)
        file_path = os.path.join(nova_path, filename)
        
        c = canvas.Canvas(file_path)
        y_position = 800
        
        # Simple word wrap mechanism to prevent text going off the side
        for line in content.split('\n'):
            words = line.split()
            current_line = ""
            for word in words:
                if len(current_line + word) > 90:
                    c.drawString(50, y_position, current_line)
                    y_position -= 15
                    current_line = word + " "
                else:
                    current_line += word + " "
            c.drawString(50, y_position, current_line)
            y_position -= 20
            
            # Start a new page if we run out of vertical space
            if y_position < 50:
                c.showPage()
                y_position = 800
                
        c.save()
        return f"Successfully created the PDF document '{filename}' in the NOVA workspace."
    except Exception as e:
        return f"Failed to create PDF '{filename}'. Error: {e}"

def create_presentation(filename: str, title: str, content: str):
    """
    Creates a new Microsoft PowerPoint presentation (.pptx) on the user's Desktop.
    Args:
        filename: The name of the presentation file to create (MUST end with .pptx).
        title: The title text for the title slide.
        content: A summary or bullet points to put on the content slide.
    """
    try:
        if not filename.endswith('.pptx'):
            filename += '.pptx'
            
        nova_path = r"d:\projects\degree\jarvis\voice_ai_assistant_free_ai\NOVA"
        os.makedirs(nova_path, exist_ok=True)
        file_path = os.path.join(nova_path, filename)
        
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title_shape.text = title
        subtitle.text = "Generated by AI Agent"
        
        # Content Slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide2.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = f"{title} Overview"
        tf = body_shape.text_frame
        tf.text = content
        
        prs.save(file_path)
        return f"Successfully created the PowerPoint presentation '{filename}' in the NOVA workspace."
    except Exception as e:
        return f"Failed to create PowerPoint '{filename}'. Error: {e}"

def check_pc_health():
    """
    Checks the current health and status of the user's computer.
    It returns the CPU usage percentage, the RAM (memory) usage percentage, and the name of the process consuming the most memory.
    Use this if the user asks why their PC is slow, or asks for a system diagnostic.
    """
    try:
        cpu_percent = psutil.cpu_percent(interval=1)
        mem = psutil.virtual_memory()
        ram_percent = mem.percent
        
        # Find the process consuming the most memory
        top_process_name = "Unknown"
        top_process_mem = 0
        
        for proc in psutil.process_iter(['name', 'memory_info']):
            try:
                mem_info = proc.info.get('memory_info')
                if mem_info and mem_info.rss > top_process_mem:
                    top_process_mem = mem_info.rss
                    top_process_name = proc.info.get('name', 'Unknown')
            except (psutil.NoSuchProcess, psutil.AccessDenied, psutil.ZombieProcess):
                pass
                
        top_mem_mb = top_process_mem / (1024 * 1024)
        
        health_report = (
            f"PC Diagnostics Report:\n"
            f"- CPU Usage: {cpu_percent}%\n"
            f"- RAM Usage: {ram_percent}% ({mem.used / (1024**3):.1f}GB / {mem.total / (1024**3):.1f}GB)\n"
            f"- Top Memory Hog: '{top_process_name}' is using {top_mem_mb:.1f} MB of RAM."
        )
        return health_report
    except Exception as e:
        return f"Failed to retrieve PC health status. Error: {e}"

def media_play_pause():
    """
    Toggles play or pause for the current system media (e.g., Spotify, Chrome, YouTube).
    """
    try:
        pyautogui.press("playpause")
        return "Successfully toggled media play/pause."
    except Exception as e:
        return f"Failed to toggle play/pause: {e}"

def media_next_track():
    """
    Skips to the next track in the current system media player.
    """
    try:
        pyautogui.press("nexttrack")
        return "Successfully skipped to the next track."
    except Exception as e:
        return f"Failed to skip track: {e}"

def media_volume_up():
    """
    Increases the master system volume of the computer.
    """
    try:
        # Press 5 times for a noticeable volume jump
        pyautogui.press("volumeup", presses=5) 
        return "Successfully increased system volume."
    except Exception as e:
        return f"Failed to increase volume: {e}"

def media_volume_down():
    """
    Decreases the master system volume of the computer.
    """
    try:
        # Press 5 times for a noticeable volume drop
        pyautogui.press("volumedown", presses=5)
        return "Successfully decreased system volume."
    except Exception as e:
        return f"Failed to decrease volume: {e}"

def set_system_volume(percentage: int):
    """
    Sets the absolute master system volume of the computer to a specific percentage.
    Args:
        percentage: The target volume level (0 to 100).
    """
    try:
        percentage = min(max(int(percentage), 0), 100)
        
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(
            IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        
        # Scalar volume ranges from 0.0 to 1.0
        scalar_vol = percentage / 100.0
        volume.SetMasterVolumeLevelScalar(scalar_vol, None)
        
        return f"Successfully set system volume to {percentage}%."
    except Exception as e:
        return f"Failed to set system volume. Error: {e}"
